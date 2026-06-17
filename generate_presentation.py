from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Theme colors
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(107, 114, 128)
LIGHT = RGBColor(243, 244, 246)


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = WHITE

    if subtitle:
        subtitle_box = slide.placeholders[1]
        subtitle_box.text = subtitle
        subtitle_box.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle_box.text_frame.paragraphs[0].font.color.rgb = WHITE


def add_content_slide(title, bullets, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY

    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p.space_after = Pt(8)

    if footer:
        slide.notes_slide.notes_text_frame.text = footer


# Slide content
slides = [
    ("Student Placement System", "A Digital Solution for Campus Placements"),
    ("Introduction", [
        "The project helps students track placement opportunities easily.",
        "It provides a centralized platform for students and admins.",
        "The system reduces manual errors and improves transparency."
    ]),
    ("Problem Statement", [
        "Students find it difficult to track company details and updates.",
        "Admins struggle to manage applications and placement schedules manually.",
        "There is a need for a structured digital system."
    ]),
    ("Objectives", [
        "Allow students to register and log in securely.",
        "Display company details and placement opportunities.",
        "Enable students to apply and monitor their application progress.",
        "Provide admins with a dashboard to manage data."
    ]),
    ("System Features", [
        "Student profile management",
        "Resume upload and viewing",
        "Company listing and application form",
        "Placement calendar for students",
        "Admin dashboard and application tracking"
    ]),
    ("Admin Module", [
        "Admin login and secure dashboard",
        "Add, edit, and delete company records",
        "Review and update student application status",
        "Manage placement drive calendar events"
    ]),
    ("Database Design", [
        "students table stores student details and login information.",
        "companies table stores company and package information.",
        "applications table links students with companies.",
        "placement_calendar stores drive dates and venues."
    ]),
    ("Technology Stack", [
        "Frontend: HTML, CSS, Jinja templates",
        "Backend: Python and Flask",
        "Database: MySQL",
        "File handling: resume uploads through Flask"
    ]),
    ("Workflow", [
        "Student registers and logs in.",
        "Student views companies and applies.",
        "Admin reviews applications and updates statuses.",
        "Students can check their dashboard and placement calendar."
    ]),
    ("Benefits", [
        "Saves time and effort for both students and admins.",
        "Improves communication and updates.",
        "Provides an organized placement management workflow.",
        "Makes the placement process more transparent."
    ]),
    ("Future Scope", [
        "Email notifications for students and admins",
        "Interview scheduling and online assessment support",
        "Analytics and reports for placements",
        "Mobile-friendly interface and better security"
    ]),
    ("Conclusion", [
        "The Student Placement System is a practical solution for campus recruitment.",
        "It simplifies management and improves student engagement.",
        "The project can be further expanded for real-world deployment."
    ]),
    ("Thank You", [
        "Thank you for your attention.",
        "Questions and discussion are welcome."
    ])
]

# Add slides
add_title_slide(slides[0][0], slides[0][1])
for title, bullets in slides[1:]:
    add_content_slide(title, bullets)

# Save file
output_path = Path('Student_Placement_System_Presentation.pptx')
prs.save(output_path)
print(f'Presentation created successfully: {output_path.resolve()}')
